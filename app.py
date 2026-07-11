"""
# =====================================
# 📚 EduDocs AI

Agente conversacional basado en RAG utilizando:

- Google Gemini
- LangChain
- ChromaDB
- Google Colab

Proyecto desarrollado para el Challenge Alura ONE.

---

# 🎯 Objetivo

Desarrollar un agente de inteligencia artificial capaz de responder preguntas
sobre documentos de una plataforma educativa utilizando la arquitectura RAG
(Retrieval-Augmented Generation).

---

# 🏗️ Arquitectura

```text
Documentos
    │
    ▼
PyPDFLoader
    │
    ▼
Document
    │
    ▼
RecursiveCharacterTextSplitter
    │
    ▼
Chunks
    │
    ▼
Google Embeddings
    │
    ▼
ChromaDB
    │
    ▼
Retriever
    │
    ▼
Gemini
    │
    ▼
Respuesta
```
"""

# =====================================
# Importación de librerías
# =====================================

# Librerías estándar de Python
import json
import os
import shutil
from pathlib import Path

# Librería utilizada para construir la interfaz web
import gradio as gr

# Librerías para procesar documentos
from openpyxl import load_workbook      # Excel (.xlsx)
from bs4 import BeautifulSoup           # HTML
from pptx import Presentation           # PowerPoint (.pptx)

# Permite acceder a la API Key almacenada en Google Colab (si está disponible)
try:
    from google.colab import userdata
except ImportError:
    userdata = None

# Clase base utilizada por LangChain para representar documentos
from langchain_core.documents import Document

# Loaders de LangChain para cargar diferentes formatos de archivos
from langchain_community.document_loaders import (
    PyPDFLoader,        # PDF
    Docx2txtLoader,     # Word (.docx)
    TextLoader,         # Texto plano y Markdown
    CSVLoader,          # Archivos CSV
)

# Divide documentos grandes en fragmentos (chunks)
from langchain_text_splitters import RecursiveCharacterTextSplitter

# Modelo utilizado para convertir texto en embeddings mediante Google Gemini
from langchain_google_genai import GoogleGenerativeAIEmbeddings

# Base de datos vectorial utilizada para almacenar y buscar embeddings
from langchain_community.vectorstores import Chroma

# =====================================
# Función para cargar documentos
# =====================================

def cargar_documentos(carpeta):

    # Convierte la ruta en un objeto Path para recorrer la carpeta
    carpeta = Path(carpeta)

    # Lista donde se almacenarán todos los documentos cargados
    documents: list[Document] = []

    # Recorre todos los archivos de la carpeta en orden alfabético
    for archivo in sorted(carpeta.iterdir()):

        try:

            # Obtiene la extensión del archivo en minúsculas
            extension = archivo.suffix.lower()

            # ---------- PDF ----------
            if extension == ".pdf":

                loader = PyPDFLoader(str(archivo))

                documents.extend(loader.load())

                print(f"✔ PDF cargado: {archivo.name}")

            # ---------- WORD ----------
            elif extension == ".docx":

                loader = Docx2txtLoader(str(archivo))

                documents.extend(loader.load())

                print(f"✔ Word cargado: {archivo.name}")

            # ---------- MARKDOWN ----------
            elif extension == ".md":

                loader = TextLoader(
                    str(archivo),
                    encoding="utf-8"
                )

                documents.extend(loader.load())

                print(f"✔ Markdown cargado: {archivo.name}")

            # ---------- CSV ----------
            elif extension == ".csv":

                loader = CSVLoader(str(archivo))

                documents.extend(loader.load())

                print(f"✔ CSV cargado: {archivo.name}")

            # ---------- JSON ----------
            elif extension == ".json":

                # Lee el contenido del archivo JSON
                with open(archivo, "r", encoding="utf-8") as f:
                    contenido = json.load(f)

                # Convierte el JSON en texto legible
                texto = json.dumps(
                    contenido,
                    indent=2,
                    ensure_ascii=False
                )

                # Crea manualmente un objeto Document
                documents.append(
                    Document(
                        page_content=texto,
                        metadata={
                            "source": str(archivo),
                            "tipo": "json"
                        }
                    )
                )

                print(f"✔ JSON cargado: {archivo.name}")

            # ---------- HTML ----------
            elif extension == ".html":

                # Lee el contenido HTML
                with open(archivo, "r", encoding="utf-8") as f:
                    html = f.read()

                # Elimina las etiquetas HTML y conserva únicamente el texto
                soup = BeautifulSoup(html, "html.parser")

                texto = soup.get_text(
                    separator="\n",
                    strip=True
                )

                # Crea un Document con el texto extraído
                documents.append(
                    Document(
                        page_content=texto,
                        metadata={
                            "source": str(archivo),
                            "tipo": "html"
                        }
                    )
                )

                print(f"✔ HTML cargado: {archivo.name}")

            # ---------- POWERPOINT ----------
            elif extension == ".pptx":

                # Abre la presentación
                presentacion = Presentation(str(archivo))

                # Recorre todas las diapositivas
                for numero_slide, slide in enumerate(presentacion.slides, start=1):

                    texto = ""

                    # Extrae el texto de cada elemento de la diapositiva
                    for shape in slide.shapes:

                        if hasattr(shape, "text"):
                            texto += shape.text + "\n"

                    texto = texto.strip()

                    # Crea un Document únicamente si la diapositiva contiene texto
                    if texto:

                        documents.append(
                            Document(
                                page_content=texto,
                                metadata={
                                    "source": str(archivo),
                                    "tipo": "powerpoint",
                                    "slide": numero_slide
                                }
                            )
                        )

                print(f"✔ PowerPoint cargado: {archivo.name}")

            # ---------- EXCEL ----------
            elif extension == ".xlsx":

                # Abre el archivo Excel
                workbook = load_workbook(archivo)

                texto = ""

                # Recorre todas las hojas del libro
                for hoja in workbook.worksheets:

                    texto += f"\nHoja: {hoja.title}\n"

                    # Convierte cada fila en texto
                    for fila in hoja.iter_rows(values_only=True):

                        fila_texto = " | ".join(
                            str(celda)
                            for celda in fila
                            if celda is not None
                        )

                        texto += fila_texto + "\n"

                # Crea un Document con todo el contenido del Excel
                documents.append(
                    Document(
                        page_content=texto,
                        metadata={
                            "source": str(archivo),
                            "tipo": "excel"
                        }
                    )
                )

                print(f"✔ Excel cargado: {archivo.name}")

            # ---------- FORMATO NO SOPORTADO ----------
            else:

                print(f"⏭ Archivo no soportado: {archivo.name}")

        # Captura cualquier error sin detener la carga del resto de documentos
        except Exception as e:

            print(f"❌ Error cargando {archivo.name}: {e}")

    # Muestra la cantidad total de documentos cargados
    print(f"\nTotal de documentos cargados: {len(documents)}")

    return documents

if not Path("documents").exists():
    raise FileNotFoundError(
        "No existe la carpeta 'documents'."
    )

# =====================================
# Cargar documentos
# =====================================

# Carga todos los documentos de la carpeta y los convierte en objetos Document
documents = cargar_documentos("documents")

if not documents:
    raise ValueError(
        "No se encontraron documentos para indexar."
    )

# =====================================
# Dividir los documentos en chunks
# =====================================

# Configura el Text Splitter que se encargará de dividir
# los documentos en fragmentos más pequeños
text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=800,      # Tamaño máximo de cada chunk
    chunk_overlap=100     # Superposición entre chunks consecutivos
)

# Divide todos los documentos cargados en chunks
chunks = text_splitter.split_documents(documents)

if not chunks:
    raise ValueError(
        "No se generaron chunks."
    )

# =====================================
# Configuración de la API Key
# =====================================

# Obtiene la API Key desde Google Colab o desde una variable de entorno
if userdata is not None:
    GOOGLE_API_KEY = userdata.get("GOOGLE_API_KEY")
else:
    GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")

# Verifica que la API Key exista
if not GOOGLE_API_KEY:
    raise ValueError(
        "No se encontró la GOOGLE_API_KEY."
    )

# =====================================
# Crear el modelo de embeddings
# =====================================

# Inicializa el modelo de embeddings de Google Gemini, el cual
# convierte texto en vectores numéricos para realizar búsquedas semánticas
embeddings = GoogleGenerativeAIEmbeddings(
    model="models/gemini-embedding-2",
    google_api_key=GOOGLE_API_KEY
)

# =====================================
# Crear o cargar la base de datos vectorial
# =====================================
# Verifica si ya existe una base vectorial almacenada localmente.
#
# - Si no existe, genera los embeddings de todos los chunks y crea
#   una nueva base de datos vectorial utilizando ChromaDB.
#
# - Si la base ya existe, la reutiliza para evitar volver a generar
#   los embeddings, reduciendo el tiempo de procesamiento y el consumo
#   de la cuota de la API de Gemini.
#
# - Si la base existe pero está vacía, elimina la carpeta y la crea
#   nuevamente para asegurar que todos los documentos queden indexados.



if not os.path.exists("chroma_db"):

    print("Creando base vectorial...")

    vector_db = Chroma.from_documents(
        documents=chunks,
        embedding=embeddings,
        persist_directory="chroma_db"
    )

else:

    vector_db = Chroma(
        persist_directory="chroma_db",
        embedding_function=embeddings
    )

    if vector_db._collection.count() == 0:

        print("La base existe pero está vacía. Recreando...")


        shutil.rmtree("chroma_db")

        vector_db = Chroma.from_documents(
            documents=chunks,
            embedding=embeddings,
            persist_directory="chroma_db"
        )

# =====================================
# Verificar la indexación de los documentos
# =====================================

print(f"Chunks indexados: {vector_db._collection.count()}")

# =====================================
# Crear el Retriever
# =====================================

# Crea el componente encargado de recuperar los fragmentos
# más relevantes desde la base de datos vectorial utilizando
# la similitud entre la pregunta y los embeddings almacenados
retriever = vector_db.as_retriever(
    search_kwargs={
        "k": 3   # Número de chunks que se recuperarán en cada búsqueda
    }
)

# =====================================
# Importar el modelo conversacional de Gemini
# =====================================

# Importa la clase utilizada para integrar un modelo
# conversacional de Google Gemini dentro de LangChain
from langchain_google_genai import ChatGoogleGenerativeAI

# =====================================
# Configurar el modelo de lenguaje (LLM)
# =====================================

# Configura el modelo de lenguaje de Google Gemini que generará
# las respuestas utilizando el contexto recuperado por el Retriever
llm = ChatGoogleGenerativeAI(
    model="gemini-3.5-flash",          # Modelo de lenguaje
    google_api_key=GOOGLE_API_KEY,     # Credenciales de acceso a la API
    temperature=0                      # Reduce la aleatoriedad para obtener respuestas más consistentes
)

# =====================================
# Importar la plantilla de Prompt
# =====================================

# Importa la clase que permite construir el Prompt que recibirá
# el modelo de lenguaje junto con el contexto recuperado
from langchain_core.prompts import ChatPromptTemplate

# =====================================
# Crear el Prompt del agente
# =====================================

# Define las instrucciones que recibirá el modelo de lenguaje.
# El Prompt establece el comportamiento del agente indicando
# cómo debe responder y qué información puede utilizar.
prompt = ChatPromptTemplate.from_template("""
Eres un asistente virtual para una plataforma educativa.

Responde únicamente utilizando la información contenida en el contexto.

Si la respuesta no se encuentra en el contexto responde exactamente:

"No encontré esa información dentro de los documentos disponibles."

Nunca inventes información.

Si la información proviene de uno o varios documentos del contexto,
menciónalos al final de la respuesta cuando sea posible.

Contexto:
{context}

Pregunta:
{question}
""")

# =====================================
# Importar componentes para la cadena RAG
# =====================================

# Importa las clases que permiten construir el flujo de trabajo
# del agente y procesar la respuesta generada por el modelo
from langchain_core.output_parsers import StrOutputParser
from langchain_core.runnables import RunnablePassthrough

# =====================================
# Construir la cadena RAG
# =====================================

# Crea el flujo de trabajo del agente RAG. La pregunta del usuario
# se utiliza para recuperar el contexto relevante desde la base de
# datos vectorial, construir el Prompt, enviarlo al modelo Gemini
# y devolver la respuesta como texto.
rag_chain = (
    {
        # Recupera los fragmentos más relevantes desde ChromaDB
        "context": retriever,

        # Pasa la pregunta del usuario al Prompt
        "question": RunnablePassthrough()
    }

    # Construye el Prompt utilizando el contexto recuperado
    | prompt

    # Envía el Prompt al modelo de lenguaje
    | llm

    # Convierte la respuesta del modelo en texto plano
    | StrOutputParser()
)

# =====================================
# Función para consultar el agente
# =====================================

historial = []

def consultar_agente(pregunta):

    if not pregunta.strip():
        return "", "\n".join(historial)

    respuesta = rag_chain.invoke(pregunta)

    historial.append(
        f"👤 Usuario:\n{pregunta}\n\n🤖 EduDocsAI:\n{respuesta}"
    )

    historial_texto = "\n\n" + ("=" * 70 + "\n\n").join(historial)

    return respuesta, historial_texto

# =====================================
# Crear la interfaz con Gradio
# =====================================



# -------------------------------------
# Funciones para los botones de feedback
# -------------------------------------

def feedback_positivo():
    return "¡Gracias por tu retroalimentación!"

def feedback_negativo():
    return "Gracias. La respuesta ayudará a mejorar el agente."

# -------------------------------------
# Interfaz principal
# -------------------------------------

with gr.Blocks(title="📚 EduDocsAI - Asistente Inteligente para Documentos Educativos") as demo:

    gr.Markdown("""
# 📚 EduDocsAI

### Asistente Virtual basado en Inteligencia Artificial

Este agente responde preguntas utilizando únicamente la información contenida
en los documentos cargados en el sistema.

Escribe una pregunta y el agente buscará la información más relevante utilizando
la base vectorial creada con LangChain, ChromaDB y Google Gemini.
""")

    # ---------------------------------
    # Pregunta del usuario
    # ---------------------------------

    pregunta = gr.Textbox(
        label="Escribe tu pregunta",
        placeholder="Ejemplo: ¿Cómo obtengo mi certificado?"
    )

    # ---------------------------------
    # Respuesta del agente
    # ---------------------------------

    respuesta = gr.Textbox(
        label="Respuesta",
        lines=10,
        interactive=False
    )

    # ---------------------------------
    # Historial de conversación
    # ---------------------------------

    historial_chat = gr.Textbox(
        label="Historial de conversación",
        lines=12,
        interactive=False
    )

    # ---------------------------------
    # Botón para consultar al agente
    # ---------------------------------

    boton = gr.Button(
        "🚀 Consultar",
        variant="primary"
    )

    boton.click(
        fn=consultar_agente,
        inputs=pregunta,
        outputs=[
            respuesta,
            historial_chat
        ]
    )

    # ---------------------------------
    # Retroalimentación del usuario
    # ---------------------------------

    gr.Markdown("## 👍 Retroalimentación")

    gr.Markdown(
        "¿La respuesta proporcionada por el agente fue útil?"
    )

    with gr.Row():

        boton_si = gr.Button("👍 Sí")

        boton_no = gr.Button("👎 No")

    feedback = gr.Textbox(
        label="Estado",
        interactive=False
    )

    boton_si.click(
        fn=feedback_positivo,
        outputs=feedback
    )

    boton_no.click(
        fn=feedback_negativo,
        outputs=feedback
    )

if __name__ == "__main__":

    demo.launch(
        server_name="0.0.0.0",
        server_port=7860
    )